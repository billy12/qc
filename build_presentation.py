from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
GREEN_DARK  = RGBColor(0x1A, 0x5C, 0x38)   # Queen Creek green
GREEN_MID   = RGBColor(0x2E, 0x7D, 0x4F)
GOLD        = RGBColor(0xD4, 0xA0, 0x17)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF4)
DARK_GRAY   = RGBColor(0x2B, 0x2B, 0x2B)
MID_GRAY    = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank


# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_text_lines(slide, lines, l, t, w, h,
                   size=16, color=DARK_GRAY, bold_first=False, spacing=1.15):
    """lines = list of strings; each becomes its own paragraph"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb

def slide_header(slide, title, subtitle=None):
    """Green bar across top with title."""
    add_rect(slide, 0, 0, 13.33, 1.35, fill_rgb=GREEN_DARK)
    add_rect(slide, 0, 1.35, 13.33, 0.08, fill_rgb=GOLD)
    add_text(slide, title, 0.35, 0.15, 12.0, 0.85,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.9, 12.0, 0.45,
                 size=14, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

def bullet_box(slide, header, bullets, l, t, w, h, header_color=GREEN_DARK):
    """Rounded card with a header and bullet list."""
    add_rect(slide, l, t, w, 0.42, fill_rgb=header_color)
    add_text(slide, header, l+0.15, t+0.06, w-0.2, 0.35,
             size=14, bold=True, color=WHITE)
    body_h = h - 0.42
    txb = slide.shapes.add_textbox(Inches(l+0.1), Inches(t+0.48),
                                   Inches(w-0.2), Inches(body_h-0.1))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_GRAY


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=GREEN_DARK)
add_rect(s1, 0, 5.5, 13.33, 2.0, fill_rgb=GREEN_MID)
add_rect(s1, 0.4, 2.85, 12.5, 0.07, fill_rgb=GOLD)

add_text(s1, "SPRING INTO QC", 0.5, 0.9, 12.3, 1.1,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "Enhancing Engagement Beyond the Egg Hunt", 0.5, 2.05, 12.3, 0.75,
         size=22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s1, "Executive Recommendation Presentation", 0.5, 3.1, 12.3, 0.55,
         size=16, bold=False, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(s1, "Town of Queen Creek  |  Parks & Recreation  |  Special Events",
         0.5, 5.7, 12.3, 0.5,
         size=14, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s1, "Recreation Coordinator Candidate  |  July 2026",
         0.5, 6.2, 12.3, 0.5,
         size=13, bold=False, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — QUEEN CREEK COMMUNITY CONTEXT
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s2, "Who We're Serving", "Queen Creek Community Profile")

stats = [
    ("95,338", "Residents (2026)"),
    ("6.2%", "Annual Growth Rate"),
    ("$141,978", "Median HH Income"),
    ("3.41", "Avg Household Size"),
    ("37 yrs", "Median Age"),
    ("90%", "Homeownership Rate"),
]
box_w, box_h = 3.8, 1.65
for i, (num, label) in enumerate(stats):
    col = i % 3
    row = i // 3
    lx = 0.35 + col * (box_w + 0.25)
    ty = 1.65 + row * (box_h + 0.25)
    add_rect(s2, lx, ty, box_w, box_h, fill_rgb=WHITE)
    add_rect(s2, lx, ty+box_h-0.08, box_w, 0.08, fill_rgb=GOLD)
    add_text(s2, num, lx+0.1, ty+0.15, box_w-0.2, 0.85,
             size=30, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(s2, label, lx+0.1, ty+0.95, box_w-0.2, 0.55,
             size=13, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text(s2,
         '★  Town Vision: "Convenience of the city, comfort of the country"  |  FY2026 Priority: Quality of Lifestyle',
         0.35, 6.75, 12.6, 0.45,
         size=12, bold=False, color=GREEN_DARK, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CURRENT STATE
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s3, "Spring Into QC — Current State", "Frontier Family Park  |  9:00 AM – 1:00 PM  |  Annual March Event")

bullet_box(s3, "What the Event Offers", [
    "Egg Hunts — free, by age group, 9:30 AM through 12:30 PM",
    "Fast & Furriest Animal Races — chicks, bunnies & tortoises, starting 10 AM",
    "5 race heats with bracket rounds (Tortoise & the Hare, Fan Favorite)",
    "Play Zone — inflatables, bungee run, obstacle course, mega slides",
    "Tot Spot — dedicated area for ages 5 and under",
    "The Great EggScape — interactive story walk",
    "Color Blast — splatter paint experience",
    "Vendors & community partners throughout",
], 0.35, 1.55, 6.1, 5.5)

bullet_box(s3, "The Gap", [
    "3,500+ attendees at the most recent Spring Into QC",
    "Families arrive with one goal: the egg hunt",
    "Once their age group's hunt ends — the visit feels complete",
    "No clear signal about what comes next",
    "Families exit before experiencing vendors, races, or other attractions",
    "Revenue and engagement opportunities go unrealized",
    "Event elements are strong individually — but not connected",
], 6.65, 1.55, 6.35, 5.5, header_color=RGBColor(0xB0, 0x3A, 0x2E))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT THE RESEARCH REVEALED
# ══════════════════════════════════════════════════════════════════════════════
s4a = prs.slides.add_slide(blank_layout)
add_rect(s4a, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s4a, "What the Research Revealed", "Confirmed Facts · Identified Gaps · Basis for Recommendation")

# Intro strip
add_rect(s4a, 0.35, 1.6, 12.6, 0.6, fill_rgb=WHITE)
add_text(s4a,
    "Spring Into QC is a well-built event with strong programming. The gaps are operational — not programmatic.",
    0.55, 1.68, 12.1, 0.48,
    size=13, bold=False, color=GREEN_DARK, align=PP_ALIGN.CENTER, italic=True)

# Three finding cards
findings = [
    (
        "FINDING 1",
        "The Raffle Loop Is Broken",
        GREEN_DARK,
        [
            "Raffle tickets are given to attendees who correctly pick animal race winners",
            "Five race heats run from 10 AM onward — tickets accumulate throughout the morning",
            "No publicly announced Grand Prize Drawing tied to those tickets",
            "Result: Tickets feel like trinkets, not incentives to stay",
            "",
            "Source: Official event descriptions, Queen Creek Parks & Recreation",
        ]
    ),
    (
        "FINDING 2",
        "Activities Run in Silos",
        RGBColor(0xB0, 0x3A, 0x2E),
        [
            "Egg hunts, animal races, Play Zone, Great EggScape, Color Blast all exist independently",
            "No event-wide narrative or flow connecting them",
            "Families who finish one activity face an unstructured decision: what next?",
            "Without a clear path forward, the default is to leave",
            "",
            "Source: Queen Creek Independent, Chamber of Commerce event listings",
        ]
    ),
    (
        "FINDING 3",
        "The Post-Hunt Moment Is Unclaimed",
        RGBColor(0x5D, 0x4E, 0x00),
        [
            "Egg hunts run 9:30 AM – 12:30 PM, staggered by age group (already well-structured)",
            "The departure problem isn't mass exodus — it's per-wave drift after each group finishes",
            "The 30 seconds after a family's hunt ends is the critical decision window",
            "Currently: no staff, no signage, no programming targeting that exact moment",
            "",
            "Source: Event schedule analysis, Town problem statement",
        ]
    ),
]

card_w = 4.0
for i, (label, title, color, bullets) in enumerate(findings):
    lx = 0.35 + i * (card_w + 0.22)
    add_rect(s4a, lx, 2.38, card_w, 0.38, fill_rgb=color)
    add_text(s4a, label, lx+0.12, 2.41, card_w-0.2, 0.32,
             size=11, bold=True, color=GOLD)
    add_rect(s4a, lx, 2.76, card_w, 0.5, fill_rgb=RGBColor(0xE8, 0xE8, 0xE8))
    add_text(s4a, title, lx+0.12, 2.8, card_w-0.2, 0.42,
             size=14, bold=True, color=DARK_GRAY)
    add_rect(s4a, lx, 3.26, card_w, 3.9, fill_rgb=WHITE)
    txb = s4a.shapes.add_textbox(Inches(lx+0.15), Inches(3.38),
                                  Inches(card_w-0.25), Inches(3.65))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        is_source = b.startswith("Source:")
        run.text = b if not b or is_source else f"• {b}"
        run.font.size = Pt(11 if not is_source else 10)
        run.font.color.rgb = MID_GRAY if is_source else DARK_GRAY
        run.font.italic = is_source
        run.font.bold = False

add_text(s4a,
    "All three findings point to the same conclusion: the event elements are strong — the connections between them are missing.",
    0.35, 7.1, 12.6, 0.32,
    size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESEARCH FOUNDATION (inserted here, referenced before recommendation)
# ══════════════════════════════════════════════════════════════════════════════
s4b = prs.slides.add_slide(blank_layout)
add_rect(s4b, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s4b, "Research Foundation", "What Was Confirmed vs. What Was Assumed")

# Confirmed column
add_rect(s4b, 0.35, 1.55, 6.0, 0.45, fill_rgb=GREEN_DARK)
add_text(s4b, "CONFIRMED  —  From Public Sources", 0.5, 1.6, 5.7, 0.37,
         size=13, bold=True, color=WHITE)
add_rect(s4b, 0.35, 2.0, 6.0, 4.4, fill_rgb=WHITE)

confirmed = [
    ("Event Basics", "Annual March event · Frontier Family Park · 9 AM – 1 PM"),
    ("Egg Hunt Schedule", "Age-group hunts staggered 9:30 AM – 12:30 PM"),
    ("Animal Races", "Fast & Furriest · chicks, bunnies, tortoises · 5 heats from 10 AM · raffle ticket mechanic"),
    ("Other Activities", "Play Zone, Tot Spot, The Great EggScape, Color Blast · all free"),
    ("Attendance & Vendors", "3,500+ attendees · vendor booths 10'x10' · fees $75–$125 per booth"),
    ("Community Demographics", "~95K residents · 6.2% annual growth · $141,978 median HH income · avg HH size ~3.4"),
    ("Town Strategic Vision", "\"Convenience of the city, comfort of the country\" · FY26 priority: Quality of Lifestyle"),
    ("Town Survey Precedent", "Town already runs resident surveys via QR code / online link (2025 Citizen Survey)"),
    ("Vendor/Sponsor Contact", "Erica Perez · Erica.Perez@QueenCreek.org · 480-358-3719"),
]
ty = 2.12
for label, detail in confirmed:
    add_text(s4b, label, 0.5, ty, 1.8, 0.28, size=10, bold=True, color=GREEN_DARK)
    add_text(s4b, detail, 2.35, ty, 3.85, 0.36, size=10, color=DARK_GRAY)
    add_rect(s4b, 0.45, ty+0.38, 5.8, 0.02, fill_rgb=LIGHT_GRAY)
    ty += 0.42

# Assumed column
add_rect(s4b, 6.65, 1.55, 6.35, 0.45, fill_rgb=RGBColor(0xB0, 0x3A, 0x2E))
add_text(s4b, "ASSUMED  —  Flagged Per Brief Instructions", 6.8, 1.6, 6.05, 0.37,
         size=13, bold=True, color=WHITE)
add_rect(s4b, 6.65, 2.0, 6.35, 4.4, fill_rgb=WHITE)

assumed = [
    ("Early Departure Rate", "Accepted as described in Town problem statement; no exit count data available"),
    ("Vendor Revenue Impact", "Assumed lower sales correlate with early departures; no revenue figures confirmed"),
    ("Raffle Prize Drawing", "No announced drawing confirmed; assumed based on absence in all event descriptions"),
    ("Staffing Model", "Redeployment feasibility assumed; actual staffing levels not publicly available"),
    ("Sponsor Interest", "Local business appetite for prize sponsorship assumed based on comparable programs"),
    ("Dwell Time Baseline", "No existing measurement; proposed as Year 1 metric to establish"),
]
ty = 2.12
for label, detail in assumed:
    add_text(s4b, label, 6.8, ty, 1.85, 0.28, size=10, bold=True, color=RGBColor(0xB0, 0x3A, 0x2E))
    add_text(s4b, detail, 8.7, ty, 4.15, 0.38, size=10, color=DARK_GRAY)
    add_rect(s4b, 6.75, ty+0.4, 6.1, 0.02, fill_rgb=LIGHT_GRAY)
    ty += 0.46

# Data Gap callout
add_rect(s4b, 0.35, 6.48, 12.6, 0.44, fill_rgb=RGBColor(0xFD, 0xF0, 0xC8))
add_text(s4b,
    "THE DATA GAP: We don't yet know the specific reasons families leave early. Beyond this proposal, "
    "a natural next step is a formalized feedback mechanism — built on the QR/online survey method "
    "the Town already uses for resident surveys.",
    0.5, 6.53, 12.3, 0.36,
    size=10, bold=True, color=RGBColor(0x7A, 0x55, 0x00), align=PP_ALIGN.CENTER, italic=True)

# Sources footer
add_rect(s4b, 0.35, 6.98, 12.6, 0.42, fill_rgb=RGBColor(0xE8, 0xF5, 0xE9))
sources_txt = (
    "Sources: queencreekaz.gov/springintoQC  ·  Queen Creek Independent  ·  Queen Creek Chamber of Commerce  ·  "
    "Hoodline (2025)  ·  worldpopulationreview.com  ·  datausa.io  ·  census.gov/quickfacts  ·  "
    "queencreekaz.gov/government (Strategic Plan)  ·  queencreekaz.gov/FrontierFamilyPark"
)
add_text(s4b, sources_txt, 0.5, 7.02, 12.3, 0.35,
         size=9, color=GREEN_DARK, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE RECOMMENDATION: OVERVIEW (3 PILLARS)
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s4, "The Recommendation", "Two Connected Pillars — Low Cost, High Impact")

add_rect(s4, 0.35, 1.6, 12.6, 0.65, fill_rgb=WHITE)
add_text(s4,
         "Spring Into QC already has the right ingredients. The opportunity is connecting them so families experience the full event — not just the egg hunt.",
         0.55, 1.68, 12.2, 0.55,
         size=14, bold=False, color=GREEN_DARK, align=PP_ALIGN.CENTER, italic=True)

overview_pillars = [
    (
        "PILLAR 1", "Redesign the Journey", GREEN_DARK,
        "Restructure the site into one guided circular path — controlled entrances, "
        "a central egg hunt, and vendors and races woven directly into the route.",
    ),
    (
        "PILLAR 2", "Gamify Retention", RGBColor(0xB0, 0x3A, 0x2E),
        "Expand the existing raffle into a layered incentive system: an early-announced "
        "grand prize, a combined punch/bingo card, and hidden golden eggs.",
    ),
]
card_w = 6.0
for i, (label, title, color, desc) in enumerate(overview_pillars):
    lx = 0.35 + i * (card_w + 0.6)
    add_rect(s4, lx, 2.5, card_w, 0.4, fill_rgb=color)
    add_text(s4, label, lx+0.15, 2.55, card_w-0.25, 0.3, size=12, bold=True, color=GOLD)
    add_rect(s4, lx, 2.9, card_w, 3.9, fill_rgb=WHITE)
    add_text(s4, title, lx+0.15, 3.05, card_w-0.3, 0.5, size=19, bold=True, color=DARK_GRAY)
    add_text(s4, desc, lx+0.15, 3.65, card_w-0.3, 3.0, size=14, color=MID_GRAY)

add_text(s4,
    "Pillar 1 gets families moving through the whole site. Pillar 2 gives them a reason to keep moving.",
    0.35, 7.0, 12.6, 0.4,
    size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PILLAR 1: EVENT LAYOUT & PARTICIPANT JOURNEY
# ══════════════════════════════════════════════════════════════════════════════
s4p1 = prs.slides.add_slide(blank_layout)
add_rect(s4p1, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s4p1, "The Recommendation — Pillar 1", "Event Layout & Participant Journey")

add_rect(s4p1, 0.35, 1.55, 12.6, 0.55, fill_rgb=WHITE)
add_text(s4p1,
    "Restructure the site as a guided loop, not a series of disconnected stops — every family's path leads past every vendor and attraction before the exit.",
    0.5, 1.6, 12.3, 0.45,
    size=12, bold=False, color=GREEN_DARK, align=PP_ALIGN.CENTER, italic=True)

# Flow diagram
journey_steps = [
    ("ENTRANCE A / B", "Controlled entry — safety & flow"),
    ("GREAT EGG WALK", "Story-walk kicks off the journey"),
    ("VENDOR ROW", "Primary vendor & partner area"),
    ("CENTRAL EGG HUNT", "Anchors the site at the center"),
    ("ANIMAL RACES", "Flows directly from the hunt"),
    ("PERIPHERY LOOP", "Play Zone, Tot Spot, Color Blast — loop to exit"),
]
n = len(journey_steps)
arrow_w = 0.28
usable = 12.6
box_w = (usable - (n - 1) * arrow_w) / n
box_h = 1.35
fy = 2.3
for i, (label, desc) in enumerate(journey_steps):
    lx = 0.35 + i * (box_w + arrow_w)
    color = GREEN_DARK if i % 2 == 0 else GREEN_MID
    add_rect(s4p1, lx, fy, box_w, box_h, fill_rgb=color)
    add_text(s4p1, label, lx+0.05, fy+0.1, box_w-0.1, 0.45,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4p1, desc, lx+0.05, fy+0.58, box_w-0.1, 0.7,
             size=8, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < n - 1:
        add_text(s4p1, "→", lx+box_w-0.02, fy+0.38, arrow_w+0.04, 0.5,
                 size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

bullet_box(s4p1, "Design Principles", [
    "Central Egg Hunt anchors the site — everything else radiates from it",
    "Entrance A & B control traffic flow and improve child safety / headcount accountability",
    "Great Egg Walk sets the tone and starts every family on the same guided path",
    "Periphery attractions (bouncy houses, Tot Spot) sit along the return route — not off to the side",
], 0.35, 3.95, 6.1, 3.1)

bullet_box(s4p1, "Why It Works", [
    "Directly solves Finding 2 — activities stop being silos and become one connected path",
    "Families physically pass every vendor and attraction before reaching the exit",
    "No new attractions needed — this is a resequencing of the existing footprint",
    "Sets up Pillar 2: the punch/bingo card gives families a reason to complete the loop, not just walk it",
], 6.65, 3.95, 6.35, 3.1, header_color=RGBColor(0xB0, 0x3A, 0x2E))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PILLAR 2: GAMIFICATION & RETENTION INCENTIVES
# ══════════════════════════════════════════════════════════════════════════════
s4p2 = prs.slides.add_slide(blank_layout)
add_rect(s4p2, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s4p2, "The Recommendation — Pillar 2", "Gamification & Retention Incentives")

add_rect(s4p2, 0.35, 1.55, 12.6, 0.55, fill_rgb=WHITE)
add_text(s4p2,
    "The raffle ticket mechanic already exists. We expand it into a layered system that rewards families for completing the full journey — not just picking a race winner.",
    0.5, 1.6, 12.3, 0.45,
    size=12, bold=False, color=GREEN_DARK, align=PP_ALIGN.CENTER, italic=True)

gam_cards = [
    ("Grand Prize Drawing", GREEN_DARK, [
        "Announced at event open (9 AM) — families know upfront it's worth staying for",
        "High-value prize sourced via local business sponsorship — zero cost to Town",
        "Drawing held at 12:15 PM from main stage",
        "Anchors the whole incentive system",
    ]),
    ("Punch + Bingo Card", GREEN_MID, [
        "One combined card, handed out at Entrance A / B",
        "Punches for periphery attractions (Play Zone, Tot Spot, Color Blast)",
        "Bingo squares for vendor & Town partner interactions",
        "Completed card = bonus raffle ticket at the drawing table",
    ]),
    ("Hidden Golden Eggs", RGBColor(0xB0, 0x3A, 0x2E), [
        "Limited golden eggs hidden across the full venue footprint",
        "Separate from the timed age-group hunts — findable all morning",
        "Each egg redeems for bonus raffle tickets",
        "Pulls families into every corner of the site, not just the center",
    ]),
    ("Race Winning Mechanic", RGBColor(0x5D, 0x4E, 0x00), [
        "Existing mechanic: correct race picks earn raffle tickets",
        "Add a bracket-champion bonus tied to the Grand Prize pool",
        "Reinforces the raffle as the event's central storyline",
        "No new operational lift — enhancement to what's already running",
    ]),
]
gcard_w, gcard_h = 6.1, 2.15
for i, (title, color, bullets) in enumerate(gam_cards):
    col = i % 2
    row = i // 2
    lx = 0.35 + col * (gcard_w + 0.2)
    ty = 2.3 + row * (gcard_h + 0.15)
    bullet_box(s4p2, title, bullets, lx, ty, gcard_w, gcard_h, header_color=color)

add_text(s4p2,
    "Together, these four layers turn a single raffle mechanic into an event-wide incentive to explore, participate, and stay through the Grand Prize Drawing.",
    0.35, 6.95, 12.6, 0.4,
    size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OPERATIONAL PLAN
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s5, "Operational Plan", "Pre-Event Through Day-Of Execution")

phases = [
    ("8 Weeks Out", GREEN_DARK, [
        "Confirm Grand Prize Drawing format & timing (12:15 PM)",
        "Recruit local business sponsor for prize package",
        "Finalize circular site layout (entrances, vendor row, hunt, loop)",
        "Design punch + bingo card artwork",
        "Brief vendor team on updated flow",
    ]),
    ("5 Weeks Out", RGBColor(0x3E, 0x8E, 0x5C), [
        "Launch event promotion (social + newsletter)",
        "Announce participating vendors & what they're offering",
        "Tease the Grand Prize — build anticipation, don't reveal it",
        "Highlight new layout, races & golden eggs in messaging",
        "Ask vendors to cross-promote on their own channels",
    ]),
    ("2 Weeks Out", GREEN_MID, [
        "Finalize prize package (rec memberships, dining, local experiences)",
        "Print punch/bingo cards & golden egg tokens",
        "Confirm staff assignments & golden-egg hide plan",
        "Distribute updated site map to vendors & partners",
        "Prep race announcement script for emcee",
    ]),
    ("Day Of", RGBColor(0x1A, 0x70, 0x48), [
        "Hunt staff briefed on handoff protocol at check-in",
        "Entrance staff distribute cards; hide golden eggs before gates open",
        "Directional signs posted at exits & periphery loop",
        "Emcee announces race schedule & prize drawing at open (9 AM)",
        "Grand Prize Drawing announced at 12:15 PM from main stage",
    ]),
    ("Post-Event", RGBColor(0x4A, 0x4A, 0x4A), [
        "Vendor satisfaction survey (5 Q's, emailed within 48 hrs)",
        "Staff debrief — what worked, what didn't",
        "Document raffle & card submission rate as retention proxy",
        "Present findings to leadership for next year's planning",
        "Future idea: explore a family feedback loop next season",
    ]),
]

n_phases = len(phases)
gap = 0.15
box_w = (12.63 - (n_phases - 1) * gap) / n_phases
for i, (phase, color, items) in enumerate(phases):
    lx = 0.35 + i * (box_w + gap)
    add_rect(s5, lx, 1.6, box_w, 0.48, fill_rgb=color)
    add_text(s5, phase, lx+0.08, 1.65, box_w-0.12, 0.38,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s5, lx, 2.08, box_w, 5.0, fill_rgb=WHITE)
    txb = s5.shapes.add_textbox(Inches(lx+0.1), Inches(2.16),
                                 Inches(box_w-0.16), Inches(4.8))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(10.5)
        run.font.color.rgb = DARK_GRAY


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STAFFING & BUDGET
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
slide_header(s6, "Staffing & Budget", "Minimal New Investment — Maximum Leverage of Existing Resources")

# Staffing box
bullet_box(s6, "Staffing — Redeployment Model", [
    "No new hires required",
    "Hunt staff transition to handoff roles as each age group finishes",
    "Entrance staff distribute punch/bingo cards at check-in (existing role)",
    "2 staff pre-hide golden eggs before gates open (~30 min task)",
    "Emcee folds race schedule & prize drawing into existing script",
    "Event supervisor oversees layout flow, card system & handoff timing",
    "Assumption: current staffing levels cover expanded entrance & hunt operations",
], 0.35, 1.6, 6.0, 4.5)

# Budget box
add_rect(s6, 6.55, 1.6, 6.45, 4.5, fill_rgb=WHITE)
add_rect(s6, 6.55, 1.6, 6.45, 0.42, fill_rgb=GREEN_DARK)
add_text(s6, "Budget — Enhancement Cost", 6.7, 1.65, 6.2, 0.35,
         size=14, bold=True, color=WHITE)

rows = [
    ("Directional signage (6–8 signs, entrances & loop)", "$150 – $350"),
    ("Punch + bingo cards (printed, ~500 qty)", "$80 – $150"),
    ("Golden egg tokens (bulk + bonus ticket inserts)", "$40 – $75"),
    ("Prize package — sponsor-funded (assumption)", "$0"),
    ("Emcee script update", "$0"),
    ("Staff overtime (if applicable)", "$0 – $420*"),
    ("TOTAL ESTIMATED COST", "$270 – $995"),
]
row_colors = [WHITE, LIGHT_GRAY, WHITE, LIGHT_GRAY, WHITE, LIGHT_GRAY, RGBColor(0xE8, 0xF5, 0xE9)]
for j, (item, cost) in enumerate(rows):
    ry = 2.15 + j * 0.38
    add_rect(s6, 6.55, ry, 6.45, 0.38, fill_rgb=row_colors[j])
    bold = (j == len(rows)-1)
    color = GREEN_DARK if bold else DARK_GRAY
    add_text(s6, item, 6.7, ry+0.06, 4.5, 0.3,
             size=11 if not bold else 13, bold=bold, color=color)
    add_text(s6, cost, 10.85, ry+0.06, 2.0, 0.3,
             size=11 if not bold else 13, bold=bold, color=color, align=PP_ALIGN.RIGHT)

add_text(s6, "* Hunt staff transitioning to handoff roles = no additional hours if within scheduled shift",
         6.55, 4.95, 6.4, 0.38, size=10, italic=True, color=MID_GRAY)

# Revenue note
add_rect(s6, 0.35, 6.25, 12.6, 0.85, fill_rgb=RGBColor(0xE8, 0xF5, 0xE9))
add_text(s6,
         "Revenue Opportunity: Sponsor branding across the prize package, signage, and punch/bingo cards"
         " valued at $500–$2,000. Enhancement is potentially cost-neutral or better.",
         0.55, 6.32, 12.2, 0.65,
         size=12, bold=False, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FINAL RECOMMENDATION
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, 13.33, 7.5, fill_rgb=GREEN_DARK)
add_rect(s7, 0, 1.45, 13.33, 0.07, fill_rgb=GOLD)
add_rect(s7, 0, 6.0, 13.33, 0.07, fill_rgb=GOLD)

add_text(s7, "Final Recommendation", 0.5, 0.2, 12.3, 1.1,
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s7,
         "Implement two targeted, connected enhancements for Spring Into QC 2027:",
         0.5, 1.7, 12.3, 0.6,
         size=16, bold=False, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

recs = [
    ("1", "Redesign the Journey", "Restructure the site into a guided circular path — controlled entrances, a central egg hunt, and periphery attractions on the return route — so every family passes every vendor and activity before the exit."),
    ("2", "Gamify Retention", "Expand the existing raffle into a layered system: an early-announced Grand Prize Drawing, a combined punch/bingo card, hidden golden eggs, and a race bracket bonus — all built on what's already running."),
]
for i, (num, title, desc) in enumerate(recs):
    lx = 0.4 + i * 6.35
    add_rect(s7, lx, 2.45, 6.1, 3.1, fill_rgb=RGBColor(0x22, 0x6B, 0x45))
    add_text(s7, num, lx+0.15, 2.55, 0.55, 0.6,
             size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s7, title, lx+0.75, 2.62, 5.2, 0.5,
             size=17, bold=True, color=WHITE)
    add_text(s7, desc, lx+0.2, 3.2, 5.75, 2.2,
             size=13, bold=False, color=LIGHT_GRAY)

add_text(s7,
    "Looking ahead: a formalized family feedback loop is a natural next step for future seasons.",
    0.4, 5.55, 12.5, 0.35,
    size=11, bold=False, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

# Metrics
add_rect(s7, 0.4, 5.9, 12.5, 0.08, fill_rgb=RGBColor(0x22, 0x6B, 0x45))
metrics = [
    "Dwell time\n(target: +45 min)",
    "Card completion &\nraffle submission rate",
    "Vendor satisfaction\nsurvey scores",
    "Sponsor revenue\ngenerated",
]
add_text(s7, "SUCCESS METRICS", 0.4, 6.1, 2.2, 0.45,
         size=11, bold=True, color=GOLD)
for i, m in enumerate(metrics):
    add_text(s7, m, 2.7 + i*2.55, 6.08, 2.4, 0.62,
             size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ── save ───────────────────────────────────────────────────────────────────
out = r"C:\Projects\QC\SpringIntoQC_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
