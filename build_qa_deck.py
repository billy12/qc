from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GREEN_DARK  = RGBColor(0x1A, 0x5C, 0x38)
GREEN_MID   = RGBColor(0x2E, 0x7D, 0x4F)
GOLD        = RGBColor(0xD4, 0xA0, 0x17)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF4)
DARK_GRAY   = RGBColor(0x2B, 0x2B, 0x2B)
MID_GRAY    = RGBColor(0x55, 0x55, 0x55)
RED_DARK    = RGBColor(0x8B, 0x1A, 0x1A)
RED_MID     = RGBColor(0xB0, 0x3A, 0x2E)
AMBER       = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, l, t, w, h, fill=None, line=None, lw=1):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill:
        sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=14, bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb

def txt_bullets(slide, lines, l, t, w, h, size=13, color=DARK_GRAY, bold_first=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold_first and i == 0


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.33, 7.5, fill=DARK_GRAY)
rect(s, 0, 3.2, 13.33, 0.08, fill=GOLD)
rect(s, 0, 5.8, 13.33, 1.7, fill=RGBColor(0x1A, 0x1A, 0x1A))

txt(s, "Q&A PREP", 0.5, 0.7, 12.3, 1.0,
    size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Anticipated Challenges & How to Answer Them", 0.5, 1.85, 12.3, 0.7,
    size=20, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
txt(s, "Spring Into QC  |  Executive Recommendation Presentation",
    0.5, 3.45, 12.3, 0.55,
    size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txt(s, "7 Challenges  •  Know Your Hardest One  •  Practice Out Loud",
    0.5, 6.1, 12.3, 0.5,
    size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# CHALLENGE SLIDES
# ══════════════════════════════════════════════════════════════════════════════

challenges = [
    {
        "num": "01",
        "difficulty": "EASY",
        "diff_color": GREEN_DARK,
        "question": '"The raffle tickets already exist — why hasn\'t that solved the problem on its own?"',
        "why_hard": "Makes it sound like you're not proposing anything new.",
        "answer": [
            "The mechanic exists but the incentive loop is incomplete.",
            "",
            "Raffle tickets without a publicized, scheduled prize drawing are just paper. Right now there's no announced moment that gives those tickets meaning or urgency.",
            "",
            "We're not adding a new program — we're completing the loop that's already been started. That's actually an efficient use of existing infrastructure.",
        ],
        "key_phrase": "\"We're not building something new — we're activating what's already there.\"",
    },
    {
        "num": "02",
        "difficulty": "EASY",
        "diff_color": GREEN_DARK,
        "question": '"How do you know families are actually leaving early? Where\'s the data?"',
        "why_hard": "You have no attendance numbers, exit counts, or dwell time data to cite.",
        "answer": [
            "Be direct: acknowledge the data gap immediately.",
            "",
            "\"This recommendation is based on the problem statement provided by the Town, and I've flagged all assumptions as such. That said, part of what I'm recommending is establishing the measurement baseline we don't have yet.\"",
            "",
            "Propose: raffle ticket submission rate, vendor feedback surveys, and basic entry/exit counts as Year 1 metrics.",
            "",
            "Turning the data gap into a measurement plan makes you look sharp, not underprepared.",
        ],
        "key_phrase": "\"We don't have the baseline yet — so let's build it. Here's how.\"",
    },
    {
        "num": "03",
        "difficulty": "EASY",
        "diff_color": GREEN_DARK,
        "question": '"What does the prize actually cost, and where does that money come from?"',
        "why_hard": "Budget questions in front of elected officials require a concrete answer.",
        "answer": [
            "The prize is sourced from local business sponsorships — not the Town's budget.",
            "",
            "Queen Creek's median household income is $142K. Families at this event respond to experiential prizes: rec center memberships, local dining packages, activity vouchers.",
            "",
            "These are easy asks for local businesses because they get branding on the raffle announcement, signage, and emcee mentions in exchange.",
            "",
            "Net Town cost for the prize: $0. Total enhancement cost: $150–$720 (signage only).",
        ],
        "key_phrase": "\"Sponsor-funded prize. Zero cost to the Town. Mutual win for a local business.\"",
    },
    {
        "num": "04",
        "difficulty": "EASY",
        "diff_color": GREEN_DARK,
        "question": '"Post-hunt handoff staff — isn\'t that an added labor cost?"',
        "why_hard": "Staffing costs are scrutinized. Sounds like you're adding headcount.",
        "answer": [
            "It's a redeployment — not a new hire.",
            "",
            "Each age-group egg hunt requires staff while it's running. The moment that hunt concludes, those staff have no hunt to manage. We redirect them to the hunt exit for 15 minutes as a handoff ambassador.",
            "",
            "No additional hours. No overtime. No new positions.",
            "",
            "We're using the natural rhythm of the event to our advantage — staff transition from hunt management to retention management at exactly the right moment.",
        ],
        "key_phrase": "\"Same staff, smarter deployment. No new cost.\"",
    },
    {
        "num": "05",
        "difficulty": "MODERATE",
        "diff_color": AMBER,
        "question": '"Kids are wound up after the egg hunt — parents want to get home before the meltdown. You can\'t fight that."',
        "why_hard": "This is behavioral and real. A panelist with kids will know this is true.",
        "answer": [
            "Acknowledge it — then flip it.",
            "",
            "\"You're right that post-hunt energy can go either way. That's exactly why the timing of the handoff matters so much.\"",
            "",
            "The first animal race heat starts at 10 AM. A child finishing the 9:30 hunt walks out and goes straight into animals racing — a new, high-energy stimulus.",
            "",
            "We're not asking families to wait around or browse passively. We're putting the next exciting thing immediately in front of them before the crash happens.",
            "",
            "The 30 seconds after the hunt is the window. The handoff exists to own that window.",
        ],
        "key_phrase": "\"We're not fighting the energy — we're redirecting it.\"",
    },
    {
        "num": "06",
        "difficulty": "HARD",
        "diff_color": RED_MID,
        "question": '"This feels incremental. Is pointing at signs and redirecting staff really a meaningful enhancement?"',
        "why_hard": "This is your hardest question. Don't get defensive. Don't apologize for simplicity.",
        "answer": [
            "Don't flinch. Answer with conviction.",
            "",
            "\"The most operationally sound enhancements often look simple in isolation. What we're describing isn't a small tweak — it's a fundamental shift in how families experience the event arc.\"",
            "",
            "Right now: families arrive, complete their goal, and leave. Their experience is a single chapter.",
            "",
            "After this: families arrive, complete their hunt, get redirected, watch races, hold raffle tickets, and stay for the drawing. Their experience has a beginning, middle, and end.",
            "",
            "\"Low cost and high impact aren't opposites. Sometimes the best recommendation is closing the gap between what you're already offering and what people are actually experiencing.\"",
        ],
        "key_phrase": "\"We're not adding complexity — we're removing friction. That's the job.\"",
    },
    {
        "num": "07",
        "difficulty": "MODERATE",
        "diff_color": AMBER,
        "question": '"How does this specifically help vendors? That\'s the revenue problem we\'re trying to solve."',
        "why_hard": "The job is special events — vendor relationships matter. Fix 1+2 doesn't directly drive booth traffic.",
        "answer": [
            "Be honest about what this does and doesn't do.",
            "",
            "\"Fix 1 and 2 solve the precondition problem — families leave before they ever have a chance to engage with vendors. You can't drive vendor traffic if your audience has gone home.\"",
            "",
            "Longer dwell time = more opportunity for natural vendor discovery.",
            "",
            "\"For more direct vendor activation, the natural Phase 2 is positioning vendor booths along the corridor between hunt zones and the race area — families physically walk through the vendor area as part of the handoff. That's a layout decision that costs nothing and compounds the impact.\"",
            "",
            "Show you're thinking in phases, not just one answer.",
        ],
        "key_phrase": "\"This solves the precondition. Phase 2 connects vendors directly to the flow.\"",
    },
]

for qa in challenges:
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)

    # Header bar
    rect(s, 0, 0, 13.33, 1.1, fill=DARK_GRAY)
    rect(s, 0, 1.1, 13.33, 0.06, fill=GOLD)

    # Challenge number
    txt(s, qa["num"], 0.25, 0.1, 0.85, 0.88,
        size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Difficulty badge
    rect(s, 1.2, 0.28, 1.5, 0.42, fill=qa["diff_color"])
    txt(s, qa["difficulty"], 1.2, 0.3, 1.5, 0.38,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Header label
    txt(s, "ANTICIPATED CHALLENGE", 2.85, 0.1, 4.0, 0.4,
        size=10, bold=True, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Left panel — the challenge
    rect(s, 0.3, 1.3, 5.8, 5.85, fill=WHITE)
    rect(s, 0.3, 1.3, 5.8, 0.45, fill=RED_DARK)
    txt(s, "THE CHALLENGE", 0.45, 1.35, 5.5, 0.38,
        size=12, bold=True, color=WHITE)

    txt(s, qa["question"], 0.45, 1.88, 5.55, 1.5,
        size=15, bold=True, color=DARK_GRAY, italic=True)

    rect(s, 0.3, 3.5, 5.8, 0.04, fill=LIGHT_GRAY)

    txt(s, "Why it's tricky:", 0.45, 3.62, 5.5, 0.38,
        size=11, bold=True, color=RED_DARK)
    txt(s, qa["why_hard"], 0.45, 4.02, 5.5, 1.5,
        size=12, color=MID_GRAY, italic=True)

    # Right panel — the answer
    rect(s, 6.45, 1.3, 6.58, 5.85, fill=WHITE)
    rect(s, 6.45, 1.3, 6.58, 0.45, fill=GREEN_DARK)
    txt(s, "YOUR ANSWER", 6.6, 1.35, 6.3, 0.38,
        size=12, bold=True, color=WHITE)

    txt_bullets(s, qa["answer"], 6.6, 1.88, 6.3, 4.0,
                size=13, color=DARK_GRAY)

    # Key phrase footer
    rect(s, 6.45, 6.6, 6.58, 0.55, fill=RGBColor(0xE8, 0xF5, 0xE9))
    txt(s, qa["key_phrase"], 6.6, 6.65, 6.3, 0.45,
        size=12, bold=True, color=GREEN_DARK, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# FINAL SLIDE — KNOW YOUR HARDEST ONE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.33, 7.5, fill=DARK_GRAY)
rect(s, 0, 1.3, 13.33, 0.07, fill=RED_MID)
rect(s, 0, 5.9, 13.33, 0.07, fill=GOLD)

txt(s, "KNOW YOUR HARDEST ONE", 0.5, 0.2, 12.3, 1.0,
    size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 0.5, 1.55, 12.33, 4.15, fill=RGBColor(0x1A, 0x1A, 0x1A))
txt(s, "Challenge 06 is the one to own:", 0.75, 1.75, 11.8, 0.5,
    size=16, bold=True, color=RED_MID)
txt(s, '"This feels incremental. Is this really a meaningful enhancement?"',
    0.75, 2.3, 11.8, 0.7,
    size=18, bold=True, color=WHITE, italic=True)
txt(s,
    "Don't apologize. Don't get defensive. Say this with conviction:\n\n"
    "\"Low cost and high impact aren't opposites. We're not adding complexity — we're removing friction. "
    "Right now families experience one chapter of the event. After this recommendation, they experience the whole story. "
    "That's the job.\"",
    0.75, 3.1, 11.8, 2.4,
    size=14, color=LIGHT_GRAY)

txt_bullets(s, [
    "EASY to defend:   01 · 02 · 03 · 04",
    "MODERATE:         05 · 07",
    "HARDEST:          06  ← practice this one out loud, multiple times",
], 0.5, 6.1, 12.3, 1.1, size=13, color=GOLD)


out = r"C:\Projects\QC\SpringIntoQC_QA_Prep.pptx"
prs.save(out)
print(f"Saved: {out}")
